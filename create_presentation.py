#!/usr/bin/env python3
"""
TAGO Leap Hackathon Pitch Deck Generator v2.0
World-class startup pitch with proper branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# =============================================================================
# BRAND DESIGN SYSTEM
# =============================================================================

# Brand Colors (from BRAND_KIT.md)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 214, 51)      # #FFD633 - Primary accent
YELLOW_DARK = RGBColor(230, 184, 0)  # #E6B800 - Hover/active

# Opacity-based grays (simulated on black background)
GRAY_50 = RGBColor(242, 242, 242)    # white @ 95%
GRAY_100 = RGBColor(217, 217, 217)   # white @ 85%
GRAY_200 = RGBColor(179, 179, 179)   # white @ 70% - Body text
GRAY_300 = RGBColor(153, 153, 153)   # white @ 60%
GRAY_400 = RGBColor(102, 102, 102)   # white @ 40% - Captions
GRAY_500 = RGBColor(51, 51, 51)      # white @ 20%
GRAY_600 = RGBColor(26, 26, 26)      # white @ 10% - Card backgrounds
GRAY_700 = RGBColor(13, 13, 13)      # white @ 5%

# Brand Fonts
FONT_PRIMARY = "Inter"
FONT_ALT = "Space Grotesk"
FONT_MONO = "Source Code Pro"

# Layout Constants
MARGIN = 0.8  # inches
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_slide_background(slide, color=BLACK):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, font_color=WHITE, font_name=FONT_PRIMARY,
                 bold=False, italic=False, align=PP_ALIGN.LEFT,
                 vertical_anchor=MSO_ANCHOR.TOP):
    """Add a styled text box"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = False

    # Set vertical alignment
    try:
        tf.anchor = vertical_anchor
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = font_name
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align

    return txBox


def add_styled_heading(slide, left, top, width, height,
                       prefix="", accent="", suffix="",
                       font_size=48, align=PP_ALIGN.LEFT):
    """Add heading with italic yellow accent word
    Example: add_styled_heading(..., prefix="", accent="Trade", suffix=" Ideas")
    Results in: "Trade" in yellow italic, " Ideas" in white
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align

    # Add prefix (white, light)
    if prefix:
        run = p.add_run()
        run.text = prefix
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
        run.font.name = FONT_PRIMARY
        run.font.bold = False

    # Add accent word (yellow, italic)
    if accent:
        run = p.add_run()
        run.text = accent
        run.font.size = Pt(font_size)
        run.font.color.rgb = YELLOW
        run.font.name = FONT_PRIMARY
        run.font.italic = True
        run.font.bold = False

    # Add suffix (white, light)
    if suffix:
        run = p.add_run()
        run.text = suffix
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
        run.font.name = FONT_PRIMARY
        run.font.bold = False

    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=18, font_color=GRAY_200):
    """Add elegant bullet list"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"→  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_PRIMARY
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return txBox


def add_card(slide, left, top, width, height, has_border=True, border_color=GRAY_600):
    """Add a card with rounded corners and subtle border"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_700

    if has_border:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    # Adjust corner radius (via XML manipulation)
    try:
        spPr = shape._sp.spPr
        prstGeom = spPr.prstGeom
        avLst = prstGeom.avLst
        if avLst is None:
            avLst = parse_xml('<a:avLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            prstGeom.append(avLst)
        # Set corner radius to ~20% for rounded-2xl effect
        gd = parse_xml('<a:gd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="adj" fmla="val 15000"/>')
        avLst.append(gd)
    except:
        pass

    return shape


def add_badge(slide, left, top, text, bg_color=YELLOW, text_color=BLACK, width=1.4):
    """Add a small badge/pill"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.anchor = MSO_ANCHOR.MIDDLE
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = text_color
    p.font.name = FONT_PRIMARY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    return shape


def add_stat_box(slide, left, top, width, height, number, label):
    """Add a stat/metric box"""
    # Card background
    card = add_card(slide, left, top, width, height, has_border=True, border_color=YELLOW)

    # Number
    add_text_box(slide, left, top + 0.2, width, 0.6, number,
                 font_size=36, font_color=YELLOW, bold=False,
                 align=PP_ALIGN.CENTER)

    # Label
    add_text_box(slide, left, top + 0.7, width, 0.4, label,
                 font_size=14, font_color=GRAY_300,
                 align=PP_ALIGN.CENTER)

    return card


def add_flow_arrow(slide, left, top):
    """Add a simple arrow for flow diagrams"""
    add_text_box(slide, left, top, 0.6, 0.5, "→",
                 font_size=28, font_color=YELLOW, align=PP_ALIGN.CENTER)


def add_flow_box(slide, left, top, width, text, is_highlight=False):
    """Add a box for flow diagrams"""
    bg_color = YELLOW if is_highlight else GRAY_700
    text_color = BLACK if is_highlight else WHITE
    border_color = YELLOW if not is_highlight else None

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color

    if border_color and not is_highlight:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    try:
        tf.anchor = MSO_ANCHOR.MIDDLE
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = text_color
    p.font.name = FONT_PRIMARY
    p.font.bold = is_highlight
    p.alignment = PP_ALIGN.CENTER

    return shape


# =============================================================================
# SLIDE CREATION
# =============================================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1)

    # Top accent line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(MARGIN), Inches(0.4),
                                    Inches(SLIDE_WIDTH - 2*MARGIN), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()

    # Main title - TAGO
    add_text_box(slide1, MARGIN, 1.8, SLIDE_WIDTH - 2*MARGIN, 1.2, "TAGO",
                 font_size=120, font_color=WHITE, bold=False,
                 align=PP_ALIGN.CENTER)

    # LEAP in yellow italic
    add_text_box(slide1, MARGIN, 3.2, SLIDE_WIDTH - 2*MARGIN, 0.9, "LEAP",
                 font_size=80, font_color=YELLOW, italic=True, bold=False,
                 align=PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide1, MARGIN, 4.4, SLIDE_WIDTH - 2*MARGIN, 0.5,
                 "Trade ideas, not tokens.",
                 font_size=24, font_color=WHITE, italic=True,
                 align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide1, MARGIN, 5.0, SLIDE_WIDTH - 2*MARGIN, 0.4,
                 "AI-Powered Narrative Trading on Hyperliquid",
                 font_size=16, font_color=GRAY_300,
                 align=PP_ALIGN.CENTER)

    # Bounty badges - centered
    badge_y = 5.8
    add_badge(slide1, 4.8, badge_y, "PEAR")
    add_badge(slide1, 6.4, badge_y, "LI.FI")
    add_badge(slide1, 8.0, badge_y, "SALT")

    # Bottom text
    add_text_box(slide1, MARGIN, 6.8, SLIDE_WIDTH - 2*MARGIN, 0.3,
                 "Hyperstack Hackathon 2025",
                 font_size=12, font_color=GRAY_400,
                 align=PP_ALIGN.CENTER)

    # Bottom accent line
    line2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(MARGIN), Inches(7.1),
                                     Inches(SLIDE_WIDTH - 2*MARGIN), Inches(0.03))
    line2.fill.solid()
    line2.fill.fore_color.rgb = YELLOW
    line2.line.fill.background()

    # =========================================================================
    # SLIDE 2: The Opportunity
    # =========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2)

    add_styled_heading(slide2, MARGIN, 0.6, 10, 0.8,
                       prefix="The ", accent="Opportunity", suffix="",
                       font_size=48)

    add_text_box(slide2, MARGIN, 1.4, 8, 0.5,
                 "Hyperliquid is the fastest-growing perps DEX. But onboarding and trading are still too complex.",
                 font_size=18, font_color=GRAY_200, italic=True)

    # Stats row
    stats = [
        ("$50B+", "Daily Volume"),
        ("150K+", "Active Traders"),
        ("5+ Steps", "To Start Trading"),
        ("0", "Narrative Trading Tools")
    ]

    x_start = 0.8
    for i, (num, label) in enumerate(stats):
        add_stat_box(slide2, x_start + i * 3.1, 2.5, 2.8, 1.2, num, label)

    # Problem bullets
    add_text_box(slide2, MARGIN, 4.2, 6, 0.5, "Current Pain Points",
                 font_size=20, font_color=YELLOW, bold=False)

    problems = [
        "Bridging from other chains requires multiple transactions",
        "Trading interfaces show tickers, not market narratives",
        "Risk management is entirely manual",
        "Automation tools require giving up custody"
    ]
    add_bullet_list(slide2, MARGIN, 4.8, 11, 2.5, problems, font_size=17)

    # =========================================================================
    # SLIDE 3: The Problem (Visual)
    # =========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3)

    add_styled_heading(slide3, MARGIN, 0.6, 10, 0.8,
                       prefix="The ", accent="Problem", suffix="",
                       font_size=48)

    # Three problem cards
    problems_data = [
        ("Complex Onboarding", "Users need 5+ steps to bridge assets and start trading on Hyperliquid", "LI.FI"),
        ("Token-Centric Trading", "Traders pick tickers, not theses. No way to express 'AI will beat ETH'", "PEAR"),
        ("No Automated Safety", "Risk management is manual. Automation means giving up custody", "SALT")
    ]

    card_width = 3.7
    card_start = 0.8
    card_gap = 0.3

    for i, (title, desc, bounty) in enumerate(problems_data):
        x = card_start + i * (card_width + card_gap)

        # Card
        add_card(slide3, x, 1.8, card_width, 4.2, has_border=True, border_color=GRAY_500)

        # Badge
        add_badge(slide3, x + 0.2, 2.0, bounty)

        # Title
        add_text_box(slide3, x + 0.2, 2.6, card_width - 0.4, 0.8, title,
                     font_size=22, font_color=WHITE, bold=False)

        # Description
        add_text_box(slide3, x + 0.2, 3.5, card_width - 0.4, 2.2, desc,
                     font_size=15, font_color=GRAY_300)

    # Bottom insight
    add_text_box(slide3, MARGIN, 6.4, SLIDE_WIDTH - 2*MARGIN, 0.5,
                 "Traders need a simpler way to onboard, express market views, and automate safely.",
                 font_size=16, font_color=GRAY_400, italic=True,
                 align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 4: The Solution
    # =========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4)

    add_styled_heading(slide4, MARGIN, 0.5, 10, 0.8,
                       prefix="TAGO ", accent="Leap", suffix="",
                       font_size=48)

    add_text_box(slide4, MARGIN, 1.2, 10, 0.4,
                 "One platform. Three powerful integrations.",
                 font_size=18, font_color=GRAY_300, italic=True)

    # Three pillars
    pillars = [
        ("PEAR", "Trade Ideas", [
            "AI-powered narrative trading",
            "Pair & basket trades",
            "Bet-slip simplicity"
        ]),
        ("LI.FI", "One-Click Onboard", [
            "Bridge from any chain",
            "Auto-deposit to account",
            "Full route transparency"
        ]),
        ("SALT", "Robo Managers", [
            "Policy-controlled accounts",
            "Automated strategies",
            "100% non-custodial"
        ])
    ]

    pillar_width = 3.5
    pillar_start = 1.2
    pillar_gap = 0.6

    for i, (badge_text, title, features) in enumerate(pillars):
        x = pillar_start + i * (pillar_width + pillar_gap)

        # Card with yellow border for emphasis
        card = add_card(slide4, x, 2.0, pillar_width, 4.0, has_border=True, border_color=YELLOW)

        # Badge
        add_badge(slide4, x + (pillar_width - 1.4) / 2, 2.2, badge_text)

        # Title
        add_text_box(slide4, x, 2.8, pillar_width, 0.5, title,
                     font_size=24, font_color=WHITE,
                     align=PP_ALIGN.CENTER)

        # Features
        for j, feature in enumerate(features):
            add_text_box(slide4, x + 0.3, 3.5 + j * 0.6, pillar_width - 0.6, 0.5,
                         f"→ {feature}",
                         font_size=14, font_color=GRAY_200)

    # Tagline
    add_text_box(slide4, MARGIN, 6.4, SLIDE_WIDTH - 2*MARGIN, 0.5,
                 '"Trade ideas, not just single tokens" — built for Hyperliquid',
                 font_size=14, font_color=YELLOW, italic=True,
                 align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 5: PEAR - Narrative Trading
    # =========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5)

    add_badge(slide5, MARGIN, 0.5, "PEAR")
    add_styled_heading(slide5, MARGIN + 1.6, 0.4, 10, 0.7,
                       prefix="", accent="Narrative", suffix=" Trading Engine",
                       font_size=40)

    add_text_box(slide5, MARGIN, 1.1, 10, 0.4,
                 '"Build tools that let people trade ideas, not just single tokens"',
                 font_size=14, font_color=GRAY_400, italic=True)

    # Left: How it works
    add_text_box(slide5, MARGIN, 1.7, 5.5, 0.4, "How It Works",
                 font_size=18, font_color=YELLOW)

    steps = [
        "Enter your market thesis in plain English",
        "AI (Claude) generates trade suggestions",
        "Select pair or basket trade structure",
        "Adjust stake, leverage, direction",
        "Execute via Pear Protocol API"
    ]

    for i, step in enumerate(steps):
        add_text_box(slide5, MARGIN + 0.4, 2.2 + i * 0.55, 5.5, 0.5,
                     f"{i+1}.  {step}",
                     font_size=14, font_color=GRAY_200)

    # Right: Key Features
    add_text_box(slide5, 7, 1.7, 5.5, 0.4, "Key Features",
                 font_size=18, font_color=YELLOW)

    features = [
        "Narrative-first UI (themes, not tickers)",
        "Pair trades: Long A / Short B",
        "Basket trades: Long group vs benchmark",
        "Bet-slip UX (stake → direction → risk)",
        "Full trade logging with metadata"
    ]

    for i, feature in enumerate(features):
        add_text_box(slide5, 7.4, 2.2 + i * 0.55, 5, 0.5,
                     f"→  {feature}",
                     font_size=14, font_color=GRAY_200)

    # Example box
    add_card(slide5, MARGIN, 5.0, SLIDE_WIDTH - 2*MARGIN, 1.8, has_border=True, border_color=YELLOW)

    add_text_box(slide5, MARGIN + 0.3, 5.2, 4, 0.3, "Example Thesis:",
                 font_size=12, font_color=YELLOW)

    add_text_box(slide5, MARGIN + 0.3, 5.5, 11, 0.5,
                 '"I believe AI tokens will outperform Ethereum over the next month"',
                 font_size=18, font_color=WHITE, italic=True)

    add_text_box(slide5, MARGIN + 0.3, 6.1, 11, 0.4,
                 "→  AI suggests: Long RENDER, TAO, FET basket  /  Short ETH benchmark",
                 font_size=15, font_color=GRAY_300)

    # =========================================================================
    # SLIDE 6: LI.FI - Onboarding
    # =========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6)

    add_badge(slide6, MARGIN, 0.5, "LI.FI")
    add_styled_heading(slide6, MARGIN + 1.6, 0.4, 10, 0.7,
                       prefix="", accent="One-Click", suffix=" Onboarding",
                       font_size=40)

    add_text_box(slide6, MARGIN, 1.1, 10, 0.4,
                 '"Bridge users from any chain into HyperEVM using LI.FI routing"',
                 font_size=14, font_color=GRAY_400, italic=True)

    # Flow diagram - cleaner horizontal layout
    flow_y = 2.0
    flow_items = [
        ("Any Chain", False),
        ("LI.FI Router", False),
        ("HyperEVM", False),
        ("Trading Account", True)
    ]

    x_pos = 0.6
    for i, (label, is_highlight) in enumerate(flow_items):
        add_flow_box(slide6, x_pos, flow_y, 2.6, label, is_highlight)
        if i < len(flow_items) - 1:
            add_flow_arrow(slide6, x_pos + 2.6, flow_y + 0.05)
        x_pos += 3.2

    # Two columns
    add_text_box(slide6, MARGIN, 3.0, 5.5, 0.4, "User Experience",
                 font_size=18, font_color=YELLOW)

    ux_items = [
        "Select origin chain + token",
        "See full route summary & ETA",
        "Track progress in real-time",
        "Automatic deposit on arrival"
    ]
    for i, item in enumerate(ux_items):
        add_text_box(slide6, MARGIN + 0.4, 3.5 + i * 0.55, 5, 0.5,
                     f"→  {item}",
                     font_size=14, font_color=GRAY_200)

    add_text_box(slide6, 7, 3.0, 5.5, 0.4, "Technical Implementation",
                 font_size=18, font_color=YELLOW)

    tech_items = [
        "Dedicated lifi-service backend",
        "GET /onboard/options",
        "POST /onboard/quote & /track",
        "Reusable deposit component"
    ]
    for i, item in enumerate(tech_items):
        add_text_box(slide6, 7.4, 3.5 + i * 0.55, 5, 0.5,
                     f"→  {item}",
                     font_size=14, font_color=GRAY_200)

    # Highlight stat
    add_stat_box(slide6, 4.5, 5.7, 4.3, 1.2, "1 Click", "From Any Chain to Trading")

    # =========================================================================
    # SLIDE 7: SALT - Robo Managers
    # =========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide7)

    add_badge(slide7, MARGIN, 0.5, "SALT")
    add_styled_heading(slide7, MARGIN + 1.6, 0.4, 10, 0.7,
                       prefix="", accent="Policy-Controlled", suffix=" Robo Managers",
                       font_size=40)

    add_text_box(slide7, MARGIN, 1.1, 10, 0.4,
                 '"Automate and manage capital without ever taking custody"',
                 font_size=14, font_color=GRAY_400, italic=True)

    # Two columns
    add_text_box(slide7, MARGIN, 1.7, 5.5, 0.4, "Policy Controls",
                 font_size=18, font_color=YELLOW)

    policies = [
        "Max Leverage: 1-10x limit",
        "Daily Notional: $100 - $1M cap",
        "Max Drawdown: 1-50% protection",
        "Allowed Pairs: Whitelist only"
    ]
    for i, item in enumerate(policies):
        add_text_box(slide7, MARGIN + 0.4, 2.2 + i * 0.55, 5, 0.5,
                     f"→  {item}",
                     font_size=14, font_color=GRAY_200)

    add_text_box(slide7, 7, 1.7, 5.5, 0.4, "Automated Strategies",
                 font_size=18, font_color=YELLOW)

    strategies = [
        "Mean Reversion: AI vs ETH",
        "SOL Ecosystem vs BTC",
        "DeFi Momentum plays",
        "60-second strategy loop"
    ]
    for i, item in enumerate(strategies):
        add_text_box(slide7, 7.4, 2.2 + i * 0.55, 5, 0.5,
                     f"→  {item}",
                     font_size=14, font_color=GRAY_200)

    # Non-custodial callout box
    add_card(slide7, MARGIN, 4.6, SLIDE_WIDTH - 2*MARGIN, 2.2, has_border=True, border_color=YELLOW)

    add_text_box(slide7, MARGIN + 0.4, 4.9, 11, 0.5,
                 "NON-CUSTODIAL BY DESIGN",
                 font_size=24, font_color=YELLOW)

    add_text_box(slide7, MARGIN + 0.4, 5.5, 11, 1,
                 "Your funds stay in YOUR policy-controlled account. The robo manager only instructs trades under strict rules — it never holds your assets. Full audit trail in strategy_runs.",
                 font_size=16, font_color=GRAY_200)

    # =========================================================================
    # SLIDE 8: Architecture
    # =========================================================================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide8)

    add_styled_heading(slide8, MARGIN, 0.5, 10, 0.7,
                       prefix="", accent="Architecture", suffix="",
                       font_size=44)

    # Service boxes
    services = [
        ("Frontend", "Next.js 14\nRainbowKit\nTailwind CSS"),
        ("pear-service", "Claude AI\nPear Protocol\nTrade Execution"),
        ("lifi-service", "LI.FI SDK\nRoute Optimization\nDeposit Flow"),
        ("salt-service", "Salt SDK\nPolicy Engine\nStrategy Loop"),
    ]

    box_width = 2.6
    box_start = 0.9
    box_gap = 0.5

    for i, (name, desc) in enumerate(services):
        x = box_start + i * (box_width + box_gap)

        add_card(slide8, x, 1.5, box_width, 2.2, has_border=True, border_color=YELLOW)

        add_text_box(slide8, x, 1.7, box_width, 0.4, name,
                     font_size=14, font_color=YELLOW,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide8, x + 0.2, 2.2, box_width - 0.4, 1.3, desc,
                     font_size=11, font_color=GRAY_200,
                     align=PP_ALIGN.CENTER)

    # External integrations
    add_text_box(slide8, MARGIN, 4.0, 10, 0.4, "External Integrations",
                 font_size=16, font_color=GRAY_400)

    externals = ["Hyperliquid", "Pear Protocol", "LI.FI", "Salt"]
    ext_start = 1.5
    ext_gap = 2.8

    for i, name in enumerate(externals):
        x = ext_start + i * ext_gap
        add_badge(slide8, x, 4.5, name, bg_color=YELLOW, text_color=BLACK, width=2.2)

    # Tech stack
    add_text_box(slide8, MARGIN, 5.4, 10, 0.3, "Tech Stack",
                 font_size=14, font_color=GRAY_400)

    add_text_box(slide8, MARGIN, 5.8, SLIDE_WIDTH - 2*MARGIN, 0.4,
                 "TypeScript  •  Turborepo  •  Fastify  •  Supabase  •  Anthropic Claude  •  Viem/Wagmi  •  pnpm",
                 font_size=14, font_color=GRAY_200,
                 align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 9: Demo Flow
    # =========================================================================
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide9)

    add_styled_heading(slide9, MARGIN, 0.5, 10, 0.7,
                       prefix="", accent="Demo", suffix=" Flow",
                       font_size=44)

    add_text_box(slide9, MARGIN, 1.1, 10, 0.4,
                 "End-to-end:  Onboard  →  Trade  →  Automate",
                 font_size=16, font_color=YELLOW, italic=True)

    # Demo steps - horizontal timeline style
    demo_steps = [
        ("1", "Onboard", "Bridge USDC via LI.FI", "LI.FI"),
        ("2", "Connect", "Wallet auto-deposits", "LI.FI"),
        ("3", "Thesis", "Enter market view", "PEAR"),
        ("4", "Trade", "AI suggests pair trade", "PEAR"),
        ("5", "Execute", "Trade via Pear API", "PEAR"),
        ("6", "Automate", "Enable robo manager", "SALT"),
    ]

    y_pos = 1.8
    row_height = 0.85

    for i, (num, title, desc, bounty) in enumerate(demo_steps):
        y = y_pos + i * row_height

        # Number circle
        circle = slide9.shapes.add_shape(MSO_SHAPE.OVAL,
                                          Inches(0.8), Inches(y),
                                          Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = YELLOW
        circle.line.fill.background()

        tf = circle.text_frame
        try:
            tf.anchor = MSO_ANCHOR.MIDDLE
        except:
            pass
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide9, 1.5, y - 0.05, 2, 0.45, title,
                     font_size=18, font_color=YELLOW)

        # Description
        add_text_box(slide9, 3.5, y, 7, 0.45, desc,
                     font_size=15, font_color=GRAY_200)

        # Badge
        add_badge(slide9, 11, y + 0.02, bounty)

    # =========================================================================
    # SLIDE 10: Bounty Checklist
    # =========================================================================
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide10)

    add_styled_heading(slide10, MARGIN, 0.5, 10, 0.7,
                       prefix="Bounty ", accent="Requirements", suffix="",
                       font_size=44)

    add_text_box(slide10, MARGIN, 1.1, 10, 0.4,
                 "All acceptance criteria satisfied",
                 font_size=16, font_color=YELLOW, italic=True)

    # Three columns
    checklists = [
        ("PEAR", [
            "✓ Trade UI starts from ideas/themes",
            "✓ Pair + basket trades wired live",
            "✓ Narrative metadata logged",
            "✓ Bet-slip UX design",
            "✓ Automation routes to service"
        ]),
        ("LI.FI", [
            "✓ lifi-service endpoints exposed",
            "✓ Composable deposit component",
            "✓ Quote, ETA, route shown",
            "✓ Progress & error states",
            "✓ Mobile responsive"
        ]),
        ("SALT", [
            "✓ Account/policy endpoints",
            "✓ Strategy loop runs periodically",
            "✓ Policy limits enforced",
            "✓ Frontend shows all state",
            "✓ Non-custodial design"
        ])
    ]

    col_width = 3.8
    col_start = 0.8
    col_gap = 0.3

    for i, (bounty, items) in enumerate(checklists):
        x = col_start + i * (col_width + col_gap)

        add_badge(slide10, x, 1.7, bounty)

        for j, item in enumerate(items):
            add_text_box(slide10, x, 2.3 + j * 0.55, col_width, 0.5, item,
                         font_size=13, font_color=GRAY_200)

    # Summary
    add_text_box(slide10, MARGIN, 5.5, SLIDE_WIDTH - 2*MARGIN, 0.6,
                 "TAGO Leap delivers on ALL THREE bounties",
                 font_size=22, font_color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide10, MARGIN, 6.1, SLIDE_WIDTH - 2*MARGIN, 0.4,
                 "with a cohesive, production-ready platform",
                 font_size=16, font_color=GRAY_300,
                 align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 11: Why TAGO Leap
    # =========================================================================
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide11)

    add_styled_heading(slide11, MARGIN, 0.5, 10, 0.7,
                       prefix="Why ", accent="TAGO Leap", suffix="",
                       font_size=44)

    differentiators = [
        ("Innovation", "First narrative trading platform on Hyperliquid"),
        ("Integration", "Seamlessly combines PEAR + LI.FI + SALT"),
        ("User Experience", "Bet-slip simplicity, not trading terminal complexity"),
        ("Safety", "Non-custodial robo managers with policy enforcement"),
        ("Accessibility", "One-click onboarding from any blockchain"),
        ("Production Ready", "Clean architecture, comprehensive documentation")
    ]

    y = 1.5
    for title, desc in differentiators:
        # Yellow marker line
        marker = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(MARGIN), Inches(y + 0.15),
                                           Inches(0.08), Inches(0.4))
        marker.fill.solid()
        marker.fill.fore_color.rgb = YELLOW
        marker.line.fill.background()

        add_text_box(slide11, MARGIN + 0.3, y, 3.5, 0.5, title,
                     font_size=20, font_color=YELLOW)
        add_text_box(slide11, 4.5, y + 0.05, 8, 0.5, desc,
                     font_size=17, font_color=GRAY_200)
        y += 0.8

    # Quote
    add_text_box(slide11, MARGIN, 6.5, SLIDE_WIDTH - 2*MARGIN, 0.5,
                 '"Trade ideas, not just single tokens" — Pear Protocol Vision',
                 font_size=14, font_color=GRAY_400, italic=True,
                 align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 12: The Ask / CTA
    # =========================================================================
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide12)

    # Top accent line
    line = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(MARGIN), Inches(0.4),
                                     Inches(SLIDE_WIDTH - 2*MARGIN), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()

    # Logo
    add_text_box(slide12, MARGIN, 1.2, SLIDE_WIDTH - 2*MARGIN, 1, "TAGO",
                 font_size=100, font_color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide12, MARGIN, 2.4, SLIDE_WIDTH - 2*MARGIN, 0.8, "LEAP",
                 font_size=70, font_color=YELLOW, italic=True,
                 align=PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide12, MARGIN, 3.5, SLIDE_WIDTH - 2*MARGIN, 0.5,
                 "The future of narrative trading is here.",
                 font_size=22, font_color=WHITE, italic=True,
                 align=PP_ALIGN.CENTER)

    # Three badges showing achievements
    badge_y = 4.4
    add_badge(slide12, 3.5, badge_y, "PEAR ✓", width=2)
    add_badge(slide12, 5.8, badge_y, "LI.FI ✓", width=2)
    add_badge(slide12, 8.1, badge_y, "SALT ✓", width=2)

    # The Ask box
    add_card(slide12, 3, 5.2, 7.333, 1.3, has_border=True, border_color=YELLOW)

    add_text_box(slide12, 3.2, 5.4, 7, 0.4, "The Ask",
                 font_size=16, font_color=YELLOW)
    add_text_box(slide12, 3.2, 5.8, 7, 0.6,
                 "Award TAGO Leap the PEAR, LI.FI, and SALT bounties",
                 font_size=18, font_color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Thank you
    add_text_box(slide12, MARGIN, 6.8, SLIDE_WIDTH - 2*MARGIN, 0.4,
                 "Thank you!",
                 font_size=28, font_color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Bottom accent line
    line2 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(MARGIN), Inches(7.1),
                                      Inches(SLIDE_WIDTH - 2*MARGIN), Inches(0.03))
    line2.fill.solid()
    line2.fill.fore_color.rgb = YELLOW
    line2.line.fill.background()

    # =========================================================================
    # Save
    # =========================================================================
    output_path = "/Users/kevinapetrei/hyperstack/TAGO_Leap_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"  • 12 slides")
    print(f"  • Inter font family")
    print(f"  • Black/White/Yellow branding")
    print(f"  • All 3 bounties covered")
    return output_path


if __name__ == "__main__":
    create_presentation()
